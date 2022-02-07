import openpyxl
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
path = "123.xlsx" #имя файла
img_path = 'НЛМК.png' # загружаем фото
prs = Presentation()

wb_obj = openpyxl.load_workbook(path)  #Открываем файл
sheet_obj = wb_obj.active #Выбираем активный лист таблицы
m_row = sheet_obj.max_row # максимальное количество столбцов

# Выводим значения в цикле
for i in range(2, m_row + 1):
    cell_obj = sheet_obj.cell(row=i, column=10) # В column= подставляем номер нужной колонки
    a=cell_obj.value
    cell_obj1 = sheet_obj.cell(row=i, column=11) # В column= подставляем номер нужной колонки
    b=cell_obj1.value
    cell_obj2 = sheet_obj.cell(row=i, column=1) # В column= подставляем номер нужной колонки
    z=cell_obj2.value
    cell_obj3 = sheet_obj.cell(row=i, column=3) # В column= подставляем номер нужной колонки
    h=cell_obj3.value
    d=str(z)+' '+str(a)
    f='Руководитель проекта: '+str(h)
    cell_obj = sheet_obj.cell(row=i, column=12) # В column= подставляем номер нужной колонки
    j=cell_obj.value
    cell_obj = sheet_obj.cell(row=i, column=13) # В column= подставляем номер нужной колонки
    k=cell_obj.value
    cell_obj = sheet_obj.cell(row=i, column=6) # В column= подставляем номер нужной колонки
    u=cell_obj.value
    
    titleLayout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(titleLayout)

    # рамещаем фото
       
    left =Inches(0.5)
    top = Inches(7)
    pic = slide.shapes.add_picture(img_path, left, top)
    
    
    #первая фигура №идей название и руководитель проекта
    t1_left = Inches(0.5)# расстояние от левого края
    t1_top = Inches(0.5)# расстояние от верха
    t1_width = Inches(9)#длинна
    t1_height = Inches(2)#ширина
    txBox1 = slide.shapes.add_textbox(Inches(0.5),Inches(0.5),Inches(9),Inches(1))
    tf1 = txBox1.text_frame.paragraphs[0]
    tf1.vertical_anchor = MSO_ANCHOR.TOP
    tf1.word_wrap = True
    tf1.margin_top = 0
    tf1.horizontal_anchor = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    run1 = tf1.add_run()
    run1.text = d+'\n'+f
    font = run1.font
    font.name = 'Calibri'
    font.size = Pt(24)#размер шрифта
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.DARK_1

   # второй текст БЫЛО/ Текущее состояние
   
    t2_left = Inches(0.5)# расстояние от левого края
    t2_top = Inches(1.5)# расстояние от верха
    t2_width = Inches(4)#длинна
    t2_height = Inches(1)#ширина
    txBox2 = slide.shapes.add_textbox(t2_left, t2_top, t2_width,t2_height)
    tf2 = txBox2.text_frame.paragraphs[0]
    tf2.horizontal_anchor = MSO_ANCHOR.TOP
    tf2.word_wrap = True
    tf2.margin_top = 0
    tf2.horizontal_anchor = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    run2 = tf2.add_run()
    run2.text = 'Было/Текущее состояние'
    tf2.word_wrap = True
    font = run2.font
    font.name = 'Calibri'
    font.size = Pt(24)#размер шрифта
    font.bold = True
    font.italic = None  
    font.color.theme_color = MSO_THEME_COLOR.DARK_1
    
    # текс описание проблемы
   
    textBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(1.6),Inches(4.5), Inches(1.0))
    textFrame = textBox3.text_frame
    textFrame.word_wrap = True
    textParagraph = textFrame.add_paragraph()
    textParagraph.text = 'Описание проблемы'
    
    #текст самой проблемы

    textBox4 = slide.shapes.add_textbox(Inches(0.8), Inches(2),Inches(4.5), Inches(4.0))
    textFrame = textBox4.text_frame
    textFrame.word_wrap = True
    textParagraph = textFrame.add_paragraph()
    textParagraph.text = b

    # возможные причины

    textBox8 = slide.shapes.add_textbox(Inches(0.5), Inches(5),Inches(4.5), Inches(4.0))
    textFrame = textBox8.text_frame
    textFrame.word_wrap = True
    textParagraph = textFrame.add_paragraph()
    textParagraph.text = 'Возможные причины'

    # текс самой причины
    """
    textBox9 = slide.shapes.add_textbox(Inches(0.8), Inches(5.5),Inches(4.5), Inches(4.0))
    textFrame = textBox9.text_frame
    textFrame.word_wrap = True
    textParagraph = textFrame.add_paragraph()
    textParagraph.text = u
    
    """
# текст Стало/ Будующее состояние
    
    t5_left = Inches(5.2)# расстояние от левого края
    t5_top = Inches(1.5)# расстояние от верха
    t5_width = Inches(4)#длинна
    t5_height = Inches(1)#ширина
    txBox5 = slide.shapes.add_textbox(t5_left, t5_top, t5_width,t5_height)
    tf5 = txBox5.text_frame.paragraphs[0]
    tf5.horizontal_anchor = MSO_ANCHOR.TOP
    tf5.word_wrap = True
    tf5.margin_top = 0
    tf5.horizontal_anchor = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    run5 = tf5.add_run()
    run5.text = 'Стало / Будущее состояние'
    tf5.word_wrap = True
    font = run5.font
    font.name = 'Calibri'
    font.size = Pt(24)#размер шрифта
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.DARK_1
    # текс описание проблемы
   
    textBox6 = slide.shapes.add_textbox(Inches(5.3), Inches(1.6),Inches(4), Inches(1.0))
    textFrame = textBox6.text_frame
    textFrame.word_wrap = True
    textParagraph = textFrame.add_paragraph()
    textParagraph.text = 'Предлагаемые изменения'

    # текст самого решения

    textBox7 = slide.shapes.add_textbox(Inches(5.6), Inches(2),Inches(4), Inches(1.0))
    textFrame = textBox7.text_frame
    textFrame.word_wrap = True
    textParagraph = textFrame.add_paragraph()
    textParagraph.text = str(j)+'\n'+'\n''Ожидаемый результат'+'\n'+'  '+str(k)
  
prs.save('123.pptx')



